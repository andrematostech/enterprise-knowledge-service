import csv
from collections.abc import Iterator
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_TEXT_TYPES = {"text/plain", "text/markdown"}
SUPPORTED_TEX_TYPES = {"text/x-tex", "application/x-tex"}
SUPPORTED_PPTX_TYPES = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}


def _extract_docx_text(file_path: Path) -> str:
    doc = Document(str(file_path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def _extract_csv_text(file_path: Path) -> str:
    rows: list[str] = []
    with file_path.open("r", encoding="utf-8", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            rows.append("\t".join(row))
    return "\n".join(rows)


def _extract_xlsx_text(file_path: Path) -> str:
    workbook = load_workbook(str(file_path), data_only=True, read_only=True)
    sheets: list[str] = []
    for sheet in workbook.worksheets:
        lines: list[str] = [f"# Sheet: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            cells = [str(value) for value in row if value is not None and str(value) != ""]
            if cells:
                lines.append("\t".join(cells))
        sheets.append("\n".join(lines))
    return "\n\n".join(sheets)


def _extract_pptx_text(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            parts.append(text)
    return "\n".join(parts)


def extract_text_from_file(path: str, content_type: str | None = None) -> str:
    file_path = Path(path)
    if not file_path.exists():
        raise ValueError("File not found")

    suffix = file_path.suffix.lower()
    if content_type == "application/pdf" or suffix == ".pdf":
        reader = PdfReader(str(file_path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if content_type in SUPPORTED_TEXT_TYPES or suffix in {".txt", ".md", ".markdown"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    if content_type in SUPPORTED_TEX_TYPES or suffix == ".tex":
        return file_path.read_text(encoding="utf-8", errors="ignore")

    if content_type in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"} or suffix == ".docx":
        return _extract_docx_text(file_path)

    if content_type in {"text/csv"} or suffix == ".csv":
        return _extract_csv_text(file_path)

    if content_type in {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"} or suffix == ".xlsx":
        return _extract_xlsx_text(file_path)

    if content_type in SUPPORTED_PPTX_TYPES or suffix == ".pptx":
        return _extract_pptx_text(file_path)

    raise ValueError("Unsupported document type")


def iter_text_chunks(file_path: Path, chunk_size: int) -> Iterator[str]:
    buffer: list[str] = []
    current_size = 0
    with file_path.open("r", encoding="utf-8", errors="ignore") as handle:
        for line in handle:
            cleaned = line.rstrip("\n")
            if not cleaned:
                continue
            buffer.append(cleaned)
            current_size += len(cleaned)
            if current_size >= chunk_size:
                yield "\n".join(buffer)
                buffer = []
                current_size = 0
    if buffer:
        yield "\n".join(buffer)


def iter_csv_chunks(file_path: Path, rows_per_chunk: int = 200, chunk_size: int = 4000) -> Iterator[str]:
    batch: list[str] = []
    current_size = 0
    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            if not row:
                continue
            line = "\t".join(row)
            batch.append(line)
            current_size += len(line)
            if len(batch) >= rows_per_chunk or current_size >= chunk_size:
                yield "\n".join(batch)
                batch = []
                current_size = 0
    if batch:
        yield "\n".join(batch)


def iter_streamable_chunks(
    path: str,
    content_type: str | None,
    chunk_size: int,
    rows_per_chunk: int = 200,
) -> Iterator[str] | None:
    file_path = Path(path)
    suffix = file_path.suffix.lower()
    if content_type == "text/csv" or suffix == ".csv":
        return iter_csv_chunks(file_path, rows_per_chunk=rows_per_chunk, chunk_size=chunk_size)
    if content_type in SUPPORTED_TEXT_TYPES or suffix == ".txt":
        return iter_text_chunks(file_path, chunk_size=chunk_size)
    return None