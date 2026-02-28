from pathlib import Path

from openpyxl import Workbook
from docx import Document
from pptx import Presentation

from app.utils.documents import extract_text_from_file


def test_extract_docx_includes_paragraph_and_table(tmp_path: Path) -> None:
    doc_path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Hello DOCX")
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Cell A1"
    doc.save(doc_path)

    text = extract_text_from_file(str(doc_path), None)

    assert "Hello DOCX" in text
    assert "Cell A1" in text


def test_extract_csv_tab_separated_rows(tmp_path: Path) -> None:
    csv_path = tmp_path / "sample.csv"
    csv_path.write_text("a,b\n1,2\n", encoding="utf-8")

    text = extract_text_from_file(str(csv_path), "text/csv")

    assert "a\tb" in text
    assert "1\t2" in text


def test_extract_xlsx_includes_sheet_headers_and_values(tmp_path: Path) -> None:
    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet_customers = workbook.active
    sheet_customers.title = "Customers"
    sheet_customers.append(["name", "email"])
    sheet_customers.append(["Alice", "alice@x.com"])

    sheet_orders = workbook.create_sheet("Orders")
    sheet_orders.append(["id", "total"])
    sheet_orders.append([1, 10.5])

    workbook.save(xlsx_path)

    text = extract_text_from_file(str(xlsx_path), None)

    assert "# Sheet: Customers" in text
    assert "name" in text
    assert "Alice" in text
    assert "# Sheet: Orders" in text
    assert "id" in text
    assert "10.5" in text


def test_extract_tex_as_text(tmp_path: Path) -> None:
    tex_path = tmp_path / "sample.tex"
    tex_path.write_text("\\section{Intro}\nHello TeX", encoding="utf-8")

    text = extract_text_from_file(str(tex_path), "text/x-tex")

    assert "Hello TeX" in text


def test_extract_pptx_includes_slide_text(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(0, 0, 300, 100)
    textbox.text_frame.text = "Hello PPTX"
    presentation.save(pptx_path)

    text = extract_text_from_file(str(pptx_path), None)

    assert "Hello PPTX" in text
